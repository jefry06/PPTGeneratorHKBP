from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re

def votum(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
    # Remove all placeholder elements
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = "VOTUM"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    return slide

def marende(prs, docx_path, occurance, delimiter):
    doc = Document(docx_path)
    count = 0
    found_second_marende = False
    collecting = False
    current_paragraph = []
    paragraphs_between = []
    paragraph_count = 0

    for paragraph in doc.paragraphs:
        paragraph_count += 1
        if "marende" in paragraph.text.lower():
            count += 1
            if count == occurance:
                found_second_marende = True
                collecting = True
                text = paragraph.text[paragraph.text.lower().find("marende"):]
                parts = text.split("“")
                before_quote = parts[0].strip()  # Text before the quote
                quoted_text, after_quote = parts[1].split("”")  # Quoted text and remaining part

                # Clean up the remaining part
                after_quote = after_quote.strip()

                # Combine the parts with newline characters
                result = f"{before_quote}\n“{quoted_text}”\n{after_quote}"
                current_paragraph.append(result)
                continue

        if collecting:
            if delimiter in paragraph.text.lower():
                break

            if "♫♪♫" in paragraph.text or "musik" in paragraph.text.lower():
                if current_paragraph:
                    paragraphs_between.append("\n".join(current_paragraph))
                    current_paragraph = []
                    cleaned_text = paragraph.text.replace("♫♪♫", "").strip()
                    if "musik" in cleaned_text.lower():
                        cleaned_text = "MUSIK"
                    current_paragraph.append(cleaned_text)
            else:
                current_paragraph.append(paragraph.text)


    if collecting and current_paragraph:
        paragraphs_between.append("\n".join(current_paragraph))

    for paragraph in paragraphs_between:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        # Remove all placeholder elements
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = paragraph
        p.alignment = PP_ALIGN.CENTER
        if len(paragraph) > 250 :
            p.font.size = Pt(36)
        elif len(paragraph) > 200 :
            p.font.size = Pt(40)
        else:
            p.font.size = Pt(48)

    return paragraph_count

def patik(prs, docx_path):
    doc = Document(docx_path)
    count = 0
    found_second_marende = False
    collecting = False
    current_paragraph = []
    paragraphs_between = []
    paragraph_count = 0
    delimiter = "marende"

    for paragraph in doc.paragraphs:
        paragraph_count += 1
        if "p a t i k".lower() in paragraph.text.lower():
            collecting = True
            paragraphs_between.append("P A T I K")
            text = paragraph.text[paragraph.text.lower().find("p a t i k".lower()):]
            text = text.replace("P a t i k :", "").strip()
            current_paragraph.append(text)
            continue

        if collecting:
            if delimiter in paragraph.text.lower():
                break

            current_paragraph.append(paragraph.text)

    if collecting and current_paragraph:
        paragraphs_between.append("\n".join(current_paragraph))


    for paragraph in paragraphs_between:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        # Remove all placeholder elements
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        if(paragraph == "P A T I K" or paragraph == "PATIK"):
            p.font.size = Pt(54)
        else:
            if len(paragraph) > 230 :
                p.font.size = Pt(14)
            elif len(paragraph) > 200 :
                p.font.size = Pt(36)
            else:
                p.font.size = Pt(48)
        p.text = paragraph
        p.alignment = PP_ALIGN.CENTER

        # print(paragraph)

    return paragraph_count

def extract_page_one_text(docx_path):
    """
    Extract all text from the first page of a DOCX file, stopping at the first page break.
    :param docx_path: Path to the DOCX file.
    :return: Text content of the first page as a single string.
    """
    doc = Document(docx_path)
    page_text = []

    # Read paragraphs until we find a page break
    for paragraph in doc.paragraphs:
        # Skip empty paragraphs
        if not paragraph.text.strip():
            continue

        # Check if this paragraph contains a page break
        has_page_break = False
        for run in paragraph.runs:
            if "w:br" in run._element.xml and 'type="page"' in run._element.xml:
                has_page_break = True
                break

        # Add text before the page break
        if not has_page_break:
            page_text.append(paragraph.text.strip())
        else:
            # If we found a page break, add the text before it and stop
            break

    return "\n".join(page_text)

def create_cover_slide(prs, text):
    """
    Create a cover slide with the given text.
    :param prs: Presentation object.
    :param text: Text content for the slide.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Remove all placeholder elements
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    text_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for line in text.split("\n"):
        p = text_frame.add_paragraph()
        match = re.search(r"topik|huria", line, re.IGNORECASE)
        # Check if "TOPIK" or "HURIA" (case-insensitive) is found
        if match:
            # Insert a blank line before the matched word
            modified_content = line[:match.start()] + "\n" + line[match.start():]
        else:
            # If neither word is found, retain the original content
            modified_content = line

        p.font.size = Pt(36)  # Customize font size
        p.text = modified_content
        p.alignment = PP_ALIGN.CENTER  # Align text to the center

        if "partording".lower() in p.text.lower():
            p.font.size = Pt(48)  # Customize font size
        elif "ev".lower() in p.text.lower() or "ep".lower() in p.text.lower():
            p.font.size = Pt(24)  # Customize font size
        elif ",".lower() in p.text.lower():
            p.font.size = Pt(24)  # Customize font size
        elif "huria".lower() in p.text.lower():
            p.font.size = Pt(28)  # Customize font size
        else:
            p.font.size = Pt(36)  # Customize font size


def add_slide(prs, docx_path, param):


    # Extract text from docx that matches param
    doc = Document(docx_path)
    found_text = []
    for paragraph in doc.paragraphs:
        if param.lower() in paragraph.text.lower():
            # Remove leading non-alphabetic characters while keeping the rest
            cleaned_text = paragraph.text.lstrip(''.join(c for c in paragraph.text if not c.isalpha() and not c.isspace()))
            # Remove multiple spaces and strip
            cleaned_text = ' '.join(cleaned_text.split())
            if cleaned_text:  # Only add if there's text after cleaning
                if cleaned_text.endswith(':'):
                    cleaned_text = cleaned_text[:-1]
                found_text.append(cleaned_text)

    if found_text:
        # Add the found text to slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        # Remove all placeholder elements
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        for text in found_text:
            p = text_frame.add_paragraph()
            p.text = text.upper()
            p.font.size = Pt(48)
            p.alignment = PP_ALIGN.CENTER

def epistel(prs, docx_path):
    doc = Document(docx_path)
    current_speaker = None
    paragraphs = []  # List to store the final grouped paragraphs
    current_paragraph = []  # Temporarily stores lines for a single speaker
    found_epistel = False
    found_marende = False

    # Process each paragraph in the document
    for para in doc.paragraphs:
        text = para.text.strip()  # Clean up extra whitespace
        if not text:  # Skip empty lines
            continue

        if "E P I S T E L".lower() in text.lower():
            found_epistel = True
            continue

        if found_epistel and "marende".lower() in text.lower():
            found_marende = True
            continue

        if found_epistel and not found_marende:
            # Check if the line starts with "U:" or "H:"
            if text.startswith("U	:") or text.startswith("H	:") or text.startswith("H :") or text.startswith("U :"):
                # Save the current paragraph if it exists
                if current_paragraph:
                    paragraphs.append("\n".join(current_paragraph))
                    current_paragraph = []  # Reset for the next speaker

                # Start a new paragraph
                current_speaker = text[:2]  # Extract the speaker identifier
                current_paragraph.append(text)
            else:
                # Add the line to the current paragraph
                current_paragraph.append(text)

    # Append the last paragraph if any
    if current_paragraph:
        paragraphs.append("\n".join(current_paragraph))

    # Create slides for each paragraph
    for paragraph in paragraphs:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Remove all placeholder elements
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.add_paragraph()
        p.text = paragraph
        if len(paragraph) > 250 :
            p.font.size = Pt(36)
        elif len(paragraph) > 200 :
            p.font.size = Pt(40)
        else:
            p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER




def convert_with_cover(docx_path, pptx_path):
    """
    Convert DOCX to PPTX with the first slide as a cover using the first page's content.
    :param docx_path: Path to the DOCX file.
    :param pptx_path: Path to save the PPTX file.
    """
    # Extract text for the cover slide
    cover_text = extract_page_one_text(docx_path)

    # Create PowerPoint presentation
    prs = Presentation()

    # Add cover slide
    create_cover_slide(prs, cover_text)
    marende(prs, docx_path, 1, "votum".lower())
    add_slide(prs, docx_path, "Votum".lower())
    marende(prs, docx_path, 2, "p a t i k".lower())
    patik(prs, docx_path)
    marende(prs, docx_path, 3, "MANOPOTI DOSA".lower())
    add_slide(prs, docx_path, "MANOPOTI DOSA".lower())
    marende(prs, docx_path, 4, "E P I S T E L".lower())
    add_slide(prs, docx_path, "E P I S T E L".lower())
    epistel(prs, docx_path)
    marende(prs, docx_path, 5, "MANGHATINDANGHON HAPORSEAON".lower())
    add_slide(prs, docx_path, "MANGHATINDANGHON HAPORSEAON".lower())
    add_slide(prs, docx_path, "Koor".lower())
    add_slide(prs, docx_path, "Tingting".lower())
    add_slide(prs, docx_path, "Sunggul".lower())
    marende(prs, docx_path, 6, "J A M I T A".lower())
    add_slide(prs, docx_path, "J A M I T A".lower())
    marende(prs, docx_path, 7, "Tangiang".lower())
    add_slide(prs, docx_path, "Tangiang Pelean".lower())
    add_slide(prs, docx_path, "PANGUJUNGI".lower())

    # Save presentation
    prs.save(pptx_path)
    print(f"Presentation saved at: {pptx_path}")